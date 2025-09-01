"""
PDF'den PPT'ye dönüştürme modülü.

Yaklaşım: Ghostscript ile PDF sayfalarını PNG'lere rasterize ederiz, sonra
python-pptx ile her sayfayı bir slayt olarak ekleriz.
"""

from __future__ import annotations

import os
from dataclasses import dataclass
from typing import List, Optional
from pathlib import Path
from datetime import datetime
import logging
import shutil
import subprocess

from pptx import Presentation
from pptx.util import Inches


logger = logging.getLogger(__name__)


@dataclass
class ConvertResult:
    output_path: str
    page_count: int


class PDFToPPTError(Exception):
    pass


class PDFToPPTConverter:
    def __init__(self, temp_dir: str, dpi: int = 150):
        self.temp_dir = temp_dir
        self.dpi = dpi
        os.makedirs(self.temp_dir, exist_ok=True)

    def _out_name(self, src: str) -> str:
        base = Path(src).stem
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        return f"{base}_slides_{ts}.pptx"

    def _render_pages_with_gs(self, src_pdf: str) -> List[str]:
        gs = shutil.which('gs')
        if not gs:
            raise PDFToPPTError("Ghostscript bulunamadı")
        pattern = os.path.join(self.temp_dir, "page_%03d.png")
        cmd = [
            gs,
            '-dNOPAUSE', '-dBATCH', '-dSAFER',
            '-sDEVICE=png16m',
            f'-r{self.dpi}',
            f'-sOutputFile={pattern}',
            src_pdf,
        ]
        try:
            subprocess.run(cmd, check=True)
        except subprocess.CalledProcessError as e:
            logger.error(f"Ghostscript render hatası: {e}")
            raise PDFToPPTError("PDF sayfaları render edilemedi")

        images = sorted([str(p) for p in Path(self.temp_dir).glob('page_*.png')])
        if not images:
            raise PDFToPPTError("Render edilmiş sayfa bulunamadı")
        return images

    def convert(self, src_pdf: str, out_path: Optional[str] = None) -> ConvertResult:
        if not os.path.exists(src_pdf):
            raise PDFToPPTError("Kaynak PDF bulunamadı")
        if out_path is None:
            out_path = os.path.join(self.temp_dir, self._out_name(src_pdf))

        images = self._render_pages_with_gs(src_pdf)

        prs = Presentation()
        # İlk görselin oranlarını slayta uygula
        from PIL import Image
        with Image.open(images[0]) as im:
            width_px, height_px = im.size
        width_in = width_px / float(self.dpi)
        height_in = height_px / float(self.dpi)
        prs.slide_width = Inches(width_in)
        prs.slide_height = Inches(height_in)

        # Boş mizanpaj
        blank_layout = prs.slide_layouts[6]
        for img in images:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)

        prs.save(out_path)
        return ConvertResult(output_path=out_path, page_count=len(images))

    def convert_combined(self, pdf_files: List[str], out_path: Optional[str] = None) -> ConvertResult:
        """Birden fazla PDF'i tek bir PowerPoint'e dönüştür"""
        if not pdf_files:
            raise PDFToPPTError("PDF dosyası listesi boş")
        
        if out_path is None:
            base_name = f"combined_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
            out_path = os.path.join(self.temp_dir, f"{base_name}.pptx")

        # İlk PDF'in boyutlarını referans al
        first_images = self._render_pages_with_gs(pdf_files[0])
        if not first_images:
            raise PDFToPPTError("İlk PDF render edilemedi")
        
        from PIL import Image
        with Image.open(first_images[0]) as im:
            width_px, height_px = im.size
        width_in = width_px / float(self.dpi)
        height_in = height_px / float(self.dpi)

        prs = Presentation()
        prs.slide_width = Inches(width_in)
        prs.slide_height = Inches(height_in)
        blank_layout = prs.slide_layouts[6]

        total_pages = 0
        
        # Her PDF için geçici klasör oluştur
        for i, pdf_file in enumerate(pdf_files):
            temp_dir = os.path.join(self.temp_dir, f"temp_pdf_{i}")
            os.makedirs(temp_dir, exist_ok=True)
            
            # PDF'i render et
            pattern = os.path.join(temp_dir, "page_%03d.png")
            gs = shutil.which('gs')
            if not gs:
                raise PDFToPPTError("Ghostscript bulunamadı")
            
            cmd = [
                gs,
                '-dNOPAUSE', '-dBATCH', '-dSAFER',
                '-sDEVICE=png16m',
                f'-r{self.dpi}',
                f'-sOutputFile={pattern}',
                pdf_file,
            ]
            try:
                subprocess.run(cmd, check=True)
            except subprocess.CalledProcessError as e:
                logger.error(f"Ghostscript render hatası: {e}")
                raise PDFToPPTError(f"PDF {os.path.basename(pdf_file)} render edilemedi")

            # Render edilen sayfaları PowerPoint'e ekle
            images = sorted([str(p) for p in Path(temp_dir).glob('page_*.png')])
            for img in images:
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)
                total_pages += 1

        prs.save(out_path)
        return ConvertResult(output_path=out_path, page_count=total_pages)


